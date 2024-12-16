import io
import tempfile
from pptx import Presentation
from pptx.dml.color import RGBColor

def copy_images_between_ppts(src_ppt_path, dest_ppt_path):
    """
    Copy all images from a source PowerPoint file to a destination PowerPoint file.
    The images retain their coordinates, slide placement, dimensions, and some styles like outlines.

    :param src_ppt_path: Path to the source PowerPoint file containing images.
    :param dest_ppt_path: Path to the destination PowerPoint file to receive images.
    """
    # Load the presentations
    src_ppt = Presentation(src_ppt_path)
    dest_ppt = Presentation(dest_ppt_path)

    # Iterate over slides and copy images
    for src_slide, dest_slide in zip(src_ppt.slides, dest_ppt.slides):
        for shape in src_slide.shapes:
            if shape.shape_type == 13:  # 13 corresponds to Picture
                # Extract image details
                image = shape.image
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height

                # Save the image blob to a temporary file
                with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as temp_image_file:
                    temp_image_file.write(image.blob)  # Write the image bytes to the temporary file
                    temp_image_path = temp_image_file.name  # Get the path of the temporary image file

                # Add image to the corresponding slide in destination
                new_picture = dest_slide.shapes.add_picture(temp_image_path, left, top, width, height)

                # Copy the outline (border) style if it exists
                if shape.line:
                    line = shape.line
                    if isinstance(line.color, RGBColor):
                        # If the line color is an RGB color, copy it directly
                        new_picture.line.color.rgb = line.color.rgb
                    elif hasattr(line.color, 'rgb'):
                        # If the line color is a scheme color, convert it to RGB
                        new_picture.line.color.rgb = line.color.rgb

                    new_picture.line.width = line.width

    # Save the modified destination presentation
    dest_ppt.save(dest_ppt_path)
    print(f"Images copied successfully to {dest_ppt_path}")


# Example usage
src_ppt = r"C:\Users\TechnoStar\Documents\macro\short_tool\1.pptx"
dest_ppt = r"C:\Users\TechnoStar\Documents\macro\short_tool\2.pptx"

copy_images_between_ppts(src_ppt, dest_ppt)
